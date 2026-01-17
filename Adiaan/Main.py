import streamlit as st
import pandas as pd
import os
import io
import requests
from datetime import datetime
from fpdf import FPDF
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt

# ================= 1. CONFIGURATION & STYLE =================
LOGO_PATH = "Adiaan/logo.png"
DATA_FILE = "dadar_final_report.txt"
COL_NAMES = ['Guyyaa', 'Maqaa_Abbaa_Dhimmaa', 'Araddaa', 'Qaxana', 'Gosa_Tajajjilaa', 'Maqaa_Ogeessa', 'Kafaltii_Taj']
BOT_TOKEN = "8357193631:AAHCuSnXzjZTQaglkmcS0gq-EvqnkIQLDBI" 
CHAT_ID = "7329587700" 
ARADDALEEN = ["01 Araddaa", "02 Araddaa", "03 Araddaa", "04 Araddaa", "Baadiyyaa/Ala"]

st.set_page_config(page_title="Deder City Land Office", page_icon="🏢", layout="wide")

st.markdown("""
    <style>
    [data-testid="stHeader"] {display: none !important;}
    .stApp { background: #f8f9fa; }
    [data-testid="stSidebar"] { background-color: #1b5e20 !important; }
    [data-testid="stSidebar"] * { color: #ffffff !important; }
    .card { 
        background: white; padding: 20px; border-radius: 12px; 
        box-shadow: 0 4px 6px rgba(0,0,0,0.1); text-align: center; 
        border-top: 8px solid #2e7d32; margin-bottom: 10px; 
    }
    .metric-value { font-size: 26px; font-weight: bold; color: #2e7d32; }
    </style>
    """, unsafe_allow_html=True)

# ================= 2. CORE FUNCTIONS =================
def load_data():
    if not os.path.exists(DATA_FILE) or os.stat(DATA_FILE).st_size == 0:
        return pd.DataFrame(columns=COL_NAMES)
    return pd.read_csv(DATA_FILE, sep="|", names=COL_NAMES, header=None, encoding='utf-8')

def save_data(df_to_save):
    df_to_save.to_csv(DATA_FILE, sep="|", index=False, header=False, encoding="utf-8")

def send_telegram_msg(msg):
    try:
        url = f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage"
        requests.post(url, data={"chat_id": CHAT_ID, "text": msg, "parse_mode": "HTML"})
    except: pass

def create_ppt_report(df, top_ogeessa):
    prs = Presentation()
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Gabaasa Sochii Hojii Lafa Magaalaa"
    slide.placeholders[1].text = f"Deder City Land Office\nGuyyaa: {datetime.now().strftime('%d/%m/%Y')}"
    # Slide 2: Waliigala
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Gabaasa Galii fi Maamiltoota"
    content = slide.placeholders[1].text_frame
    content.text = f"• Waliigala Galii: {df['Kafaltii_Taj'].sum():,.2f} ETB"
    content.add_paragraph().text = f"• Baay'ina Maamiltootaa: {len(df)}"
    content.add_paragraph().text = f"• Ogeessa Filatamaa: {top_ogeessa}"
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()

# ================= 3. MAIN APP =================
if 'logged_in' not in st.session_state: st.session_state.logged_in = False

if not st.session_state.logged_in:
    _, col_mid, _ = st.columns([1, 1, 1])
    with col_mid:
        st.subheader("🏢 Deder City Land Login")
        with st.form("login"):
            u = st.text_input("Username")
            p = st.text_input("Password", type="password")
            if st.form_submit_button("Seeni"):
                if u == "DAD" and p == "2026":
                    st.session_state.logged_in = True; st.rerun()
                else: st.error("Maqaa ykn Password dogoggora!")
else:
    df = load_data()
    menu = st.sidebar.radio("FILANNOO", ["📊 Dashboard", "📝 Galmee Haaraa", "📈 Gabaasa Bal'aa", "🏆 Badhaasa Ogeeyyii", "🔍 Barbaadi/Edit"])

    if menu == "📊 Dashboard":
        st.title("📊 Dashboard")
        if not df.empty:
            c1, c2, c3 = st.columns(3)
            with c1: st.markdown(f"<div class='card'><p>💰 Galii</p><p class='metric-value'>{df['Kafaltii_Taj'].sum():,.2f}</p></div>", unsafe_allow_html=True)
            with c2: st.markdown(f"<div class='card'><p>👥 Maamiltoota</p><p class='metric-value'>{len(df)}</p></div>", unsafe_allow_html=True)
            with c3: st.markdown(f"<div class='card'><p>👷 Ogeeyyii</p><p class='metric-value'>{df['Maqaa_Ogeessa'].nunique()}</p></div>", unsafe_allow_html=True)
            
            st.plotly_chart(px.area(df.groupby('Guyyaa')['Kafaltii_Taj'].sum().reset_index(), x='Guyyaa', y='Kafaltii_Taj', title="Trendii Galii"), use_container_width=True)
            
            st.write("---")
            st.subheader("🖥️ Gabaasaalee Buufadhu")
            col_ex, col_ppt = st.columns(2)
            
            # Excel Download
            output = io.BytesIO()
            with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False)
            col_ex.download_button("📥 Excel Report", output.getvalue(), f"Gabaasa_Dadar_{datetime.now().strftime('%d_%m')}.xlsx")
            
            # PPT Download
            top_man = df['Maqaa_Ogeessa'].value_counts().idxmax()
            ppt_data = create_ppt_report(df, top_man)
            col_ppt.download_button("📥 PowerPoint (PPT)", ppt_data, f"Gabaasa_Dadar_{datetime.now().strftime('%d_%m')}.pptx")
        else: st.info("Data'n galmaa'e hin jiru.")

    elif menu == "📝 Galmee Haaraa":
        st.header("📝 Galmee Tajaajilaa Haaraa")
        with st.form("reg_form", clear_on_submit=True):
            col1, col2 = st.columns(2)
            maqaa = col1.text_input("Maqaa Abbaa Dhimmaa *")
            ara = col2.selectbox("Araddaa Select *", ARADDALEEN)
            qax = col1.text_input("Qaxana / Zone")
            oge = col2.text_input("Maqaa Ogeessaa *")
            taj = st.selectbox("Gosa Tajaajilaa", ["Gibira Baaxii Gooroo", "Kaartaa Haaraa", "Jijjiirraa Maqaa", "Liizii Waggaa"])
            kaf = st.number_input("Kafaltii (ETB)", min_value=0.0)
            if st.form_submit_button("💾 Galmeessi"):
                if maqaa and oge:
                    new_row = [datetime.now().strftime('%d/%m/%Y'), maqaa, ara, qax, taj, oge, kaf]
                    df = pd.concat([df, pd.DataFrame([new_row], columns=COL_NAMES)], ignore_index=True)
                    save_data(df)
                    send_telegram_msg(f"🔔 <b>GALMEE:</b> {maqaa}\n📍 {ara} | 💰 {kaf} ETB")
                    st.success("Galmeeffameera!"); st.balloons()
                else: st.error("Maaloo, ragaa guuti.")

    elif menu == "🔍 Barbaadi/Edit":
        st.header("🔍 Barbaadi fi Sirreessi")
        q = st.text_input("Maqaa maamilaa barreessi...")
        if q:
            res = df[df['Maqaa_Abbaa_Dhimmaa'].str.contains(q, case=False, na=False)]
            for idx, row in res.iterrows():
                with st.expander(f"📄 {row['Maqaa_Abbaa_Dhimmaa']} - {row['Guyyaa']}"):
                    u_maqaa = st.text_input("Maqaa", row['Maqaa_Abbaa_Dhimmaa'], key=f"m_{idx}")
                    u_kaf = st.number_input("Kafaltii", float(row['Kafaltii_Taj']), key=f"k_{idx}")
                    if st.button("💾 Update", key=f"up_{idx}"):
                        df.at[idx, 'Maqaa_Abbaa_Dhimmaa'] = u_maqaa
                        df.at[idx, 'Kafaltii_Taj'] = u_kaf
                        save_data(df); st.success("Sirreeffameera!"); st.rerun()

    if st.sidebar.button("🚪 Ba'i"):
        st.session_state.logged_in = False; st.rerun()
